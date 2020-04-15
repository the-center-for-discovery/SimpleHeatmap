
# IMPORTING NEEDED PYTHON PACKAGES
from PIL import Image, ImageDraw, ImageFont
import matplotlib.pyplot as plt
import matplotlib.image as mpimg
import pandas as pd
from docx import Document
from docx.shared import Inches
from pptx import Presentation
from pptx.util import Inches

#*****************************SETUP******************************

images = ["southcamp.jpg","northcamp.png","ridge.png","lifecenter.png","bbhville.png","therest.png"]

df = pd.read_excel ('METADATA.xlsx')

overlay_drawing = []
loc_image = []

risk_high = []
risk_med_br = []
risk_med_li = []
risk_low_br = []
risk_low_li = []
no_risk = []

posSta = []

loc_risk =[]
out = []

data = []

# SETTING COLOR AND TRANSPARENCY
colorRed = (255, 0, 0)
colorYel = (255, 255, 0) 
colorLitYel = (255, 255, 0) 
colorOng = (255, 144, 0)
colorLitOng = (255,144,0)
colorBlnk = (255, 0, 0)

# HOW TRANSPARENT; 1 MEANING FULLY, 0 MEANING NO TRANSPARENCY

degree_transparency_red = .6
degree_transparency_yel = .6
degree_transparency_lityel = .2
degree_transparency_ong = .6
degree_transparency_litong = .2
degree_transparency_blnk = .0


# DETERMINING DEGREE OF TRANSPARENCY MASK AND ADDING ALPHA CHANNEL TO SELECTED COLOR

opacity_red = int(255 * degree_transparency_red) 
colorRed = colorRed + (opacity_red,)

opacity_yel = int(255 * degree_transparency_yel)
colorYel = colorYel + (opacity_yel,) 
opacity_lityel = int(255 * degree_transparency_lityel) 
colorLitYel = colorLitYel + (opacity_lityel,)

opacity_ong = int(255 * degree_transparency_yel) 
colorOng = colorOng + (opacity_ong,) 
opacity_litong = int(255 * degree_transparency_litong) 
colorLitOng = colorLitOng + (opacity_litong,)

opacity_blnk = int(255 * degree_transparency_blnk) 
colorBlnk = colorBlnk + (opacity_blnk,)

# LABELING COORDINATES OF WHERE BUILDINGS ARE LOCATED
sc_coords = {"Hamilton":(700,1000, 800,1100),"Berman":(980,1130,1080,1230),
               "Eichenauer":(1080, 1150, 1180, 1250), "Birch":(980,870,1080,970),
               "Willow":(1000,830,1100,930),"Smith":(1070,800,1170,900),
               "Benson":(1120,870,1220,970),"Forman":(1280, 950, 1380, 1050),
               "Oak":(1700,1150,1800,1250),"Mulberry":(1790,1220,1890,1320),
            "Sunset North":(1250,0,1400,150),"Sunset South":(1250,0,1400,150),
            "Otis Armstrong":(670,500,770,600),"Four Seasons":(120,830,220,930), 
             "Pine":(1680,1130,1780,1230)}

nc_coords = {"Thyme":(870,220,970,320),"Parsley":(860,270,960,370),
             "Rosemary ":(920,250,1020,350), "Sage":(900,300,1000,400),
            "Rapp R":(1980,750,2080,850), "Rapp West":(1950,890,2050,990)}

ri_coords = {"Ashwood":(300,130,400,230), "Acorn":(400,400,500,500), 
             "Aspen":(220,350,320,450) ,"Balsam ":(280,220,380,320), 
             "Beechnut":(380,330,480,430), "Briarwood":(150,350,250,450), 
             "Cedar":(220,170,320,270), "Chestnut":(300,400,400,500), 
             "Cypress":(220,280,320,380)}

lc_coords = {"Elm":(360,670,460,770),"Evergreen":(250,640,350,740),
             "Redwood":(410,800,510,900),"Spruce":(260,580,360,680),
             "Tulip":(320,760,420,860), "Harvest":(10,190,110,290)}

bh_coords = {"Granite":(520,190,620,290), "Slate":(450,260,550,360), 
             "Stonewall":(110,150,210,250), "Barefoot":(1120,200,1220,300),
            "Main":(1060,300,1160,400), "Railroad":(1090,465,1190,565)}

tr_coords = {"Vista":(330,290,430,390),"Wawanda":(1180,450,1280,550),
            "Sweet Hill":(1510,480,1610,580)}

##FOR LOOP TO ITERATE THROUGH EXCEL DATA***************************************************
for index, row in df.iterrows():

# 
    house = str(row['LOCATION'])
    
    pos_sta = float(row['ACTIVE'])
    pos_res = float(row['RES, POS, ACTIVE'])
    
    pend_res = float(row['RES, PEND, ACTIVE'])
    symp_res = float(row['RES, S/S, ACTIVE'])
    
    pend_sta = float(row['staff tested, pending, active'])
    staff_sym = float(row['ACTIVE S&S'])
    
    
#HIGH RISK
    if pos_sta > 0:
        risk_high.append(str(house))
    elif pos_res > 0:
        risk_high.append(str(house))
        #print(house)    
#MEDIUM RISK BRIGHT
    elif pend_res > 0:
        risk_med_br.append(str(house))
#MEDIUM RISK LIGHT
    elif symp_res > 0:
        risk_med_li.append(str(house))
        
#LOW RISK BRIGHT
    elif pend_sta > 0:
        risk_low_br.append(str(house))
#LOW RISK LIGHT
    if staff_sym > 1:
        risk_low_li.append(str(house))
    elif staff_sym <= 1:
        no_risk.append(str(house))


#******************************************************************************END

for i in images:
    index = images.index(i)
    loc_image = Image.open(images[index]).convert("RGBA")
    if i == "southcamp.jpg":
        #print (coords)
        coords = sc_coords
    elif i == "northcamp.png":
        #print (coords)
        coords = nc_coords
    elif i == "ridge.png":
        #print (coords)
        coords = ri_coords
    elif i == "lifecenter.png":
        #print (coords)
        coords = lc_coords
    elif i == "bbhville.png":
        #print (coords)
        coords = bh_coords
    elif i == "therest.png":
        #print (coords)
        coords = tr_coords

    overlay = Image.new("RGBA", loc_image.size)
    overlay_drawing = ImageDraw.Draw(overlay)

    for key, value in coords.items():
        if key in risk_high:
            data.append(key)
            overlay_drawing.ellipse(coords[key], fill=colorRed)
            loc_risk = Image.alpha_composite(loc_image, overlay)
        elif key in risk_med_br: 
            overlay_drawing.ellipse(coords[key], fill=colorOng)
            loc_risk = Image.alpha_composite(loc_image, overlay)
        elif key in risk_med_li: 
            overlay_drawing.ellipse(coords[key], fill=colorLitOng)
            loc_risk = Image.alpha_composite(loc_image, overlay)
        elif key in risk_low_br:
            overlay_drawing.ellipse(coords[key], fill=colorYel)
            loc_risk = Image.alpha_composite(loc_image, overlay)
        elif key in risk_low_li:
            overlay_drawing.ellipse(coords[key], fill=colorLitYel)
            loc_risk = Image.alpha_composite(loc_image, overlay)
        elif key in no_risk:
            loc_risk = Image.alpha_composite(loc_image, overlay)
    
    loc_risk.save('{}_risk.png'.format(i[0:len(i)-4]))
    loc_risk.show()
    
    
#***************************GENERATE WORD REPORT********************************

document = Document()

sections = document.sections
for section in sections:
    section.top_margin = Inches(1.2)
    section.bottom_margin = Inches(0)
    section.left_margin = Inches(0.25)
    section.right_margin = Inches(0.25)

#*********************************PAGE 1*****************************************
document.add_heading('SOUTH CAMPUS', 0)

document.add_picture('southcamp_risk.png', width=Inches(8))
document.add_picture('legend.png', width=Inches(8))

document.add_page_break()

#*********************************PAGE 2*****************************************
document.add_heading('NORTH CAMPUS', 0)

document.add_picture('northcamp_risk.png', width=Inches(8))
document.add_picture('legend.png', width=Inches(8))

document.add_page_break()

#****************** ***************PAGE 3*****************************************
document.add_heading('THE RIDGE', 0)

document.add_picture('ridge_risk.png', width=Inches(8))
document.add_picture('legend.png', width=Inches(8))

document.add_page_break()

#*********************************PAGE 4*****************************************
document.add_heading('THE LIFE CENTER', 0)

document.add_picture('lifecenter_risk.png', width=Inches(7))
document.add_picture('legend.png', width=Inches(7))

document.add_page_break()

#*********************************PAGE 5*****************************************
document.add_heading('MITEER, HURLEYVILLE, OUTLIERS', 0)

document.add_picture('bbhville_risk.png', width=Inches(8))
document.add_picture('therest_risk.png', width=Inches(8))
document.add_picture('legend.png', width=Inches(8))

document.save('CV19 AUTO HEATMAP REPORT.docx')


#*****************************GENERATE PDF REPORT***********************************

prs = Presentation()

blank_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_slide_layout)
report = slide.shapes

img_path = 'Legend.png'
left = Inches(1.8)
top  = Inches(6.4)
height = Inches(0.8)
pic = slide.shapes.add_picture(img_path, left, top, height=height)

img_path = 'southcamp_risk.png'
left = Inches(0.5)
top  = Inches(0.2)
height = Inches(3)
pic = slide.shapes.add_picture(img_path, left, top, height=height)

img_path = 'northcamp_risk.png'
left = Inches(5)
top  = Inches(0.2)
height = Inches(2.5)
pic = slide.shapes.add_picture(img_path, left, top, height=height)

img_path = 'ridge_risk.png'
left = Inches(0.2)
top  = Inches(3.2)
height = Inches(2.5)
pic = slide.shapes.add_picture(img_path, left, top, height=height)

img_path = 'lifecenter_risk.png'
left = Inches(3.9)
top  = Inches(3.5)
height = Inches(2.5)
pic = slide.shapes.add_picture(img_path, left, top, height=height)

img_path = 'bbhville_risk.png'
left = Inches(6.4)
top  = Inches(2.9)
height = Inches(1.5)
pic = slide.shapes.add_picture(img_path, left, top, height=height)

img_path = 'therest_risk.png'
left = Inches(6.3)
top  = Inches(4.5)
height = Inches(1.5)
pic = slide.shapes.add_picture(img_path, left, top, height=height)

prs.save('CV19 AUTO HEATMAP.pptx')